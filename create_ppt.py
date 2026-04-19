"""Generate GrabHack 2.0 submission PPT for NewAge Innovators."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import pptx.oxml.ns as nsmap
from lxml import etree
import copy

# ── Brand colours ──────────────────────────────────────────────────────────────
LIME        = RGBColor(0xC8, 0xF0, 0x00)   # #C8F000  – main bg
TEAL_DARK   = RGBColor(0x00, 0x6E, 0x5E)   # #006E5E  – headings / GrabHack text
TEAL_MID    = RGBColor(0x00, 0xA8, 0x9C)   # #00A89C  – accent / circles
CYAN_LIGHT  = RGBColor(0x00, 0xD4, 0xD0)   # #00D4D0  – highlight bar
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
BLACK       = RGBColor(0x1A, 0x1A, 0x1A)
GRAB_GREEN  = RGBColor(0x00, 0xB1, 0x4F)   # Grab logo green
GRAY_BG     = RGBColor(0xF5, 0xF5, 0xF5)   # slide body bg

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

# ── Helpers ────────────────────────────────────────────────────────────────────
def _solid_fill(shape, rgb: RGBColor):
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb

def add_rect(slide, l, t, w, h, rgb, radius=0):
    shape = slide.shapes.add_shape(
        pptx.enum.shapes.MSO_SHAPE_TYPE.AUTO_SHAPE if False else 1,  # MSO_CONNECTOR_TYPE unused
        Inches(l), Inches(t), Inches(w), Inches(h)
    )
    _solid_fill(shape, rgb)
    shape.line.fill.background()
    return shape

def add_textbox(slide, text, l, t, w, h, font_name="Calibri", font_size=18,
                bold=False, color=BLACK, align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txb

def add_header_bar(slide, title_right="GrabHack 2.0"):
    """Thin top bar: Grab (green) left, GrabHack 2.0 right."""
    bar = add_rect(slide, 0, 0, 13.33, 0.45, WHITE)
    # bottom border line
    border = add_rect(slide, 0, 0.43, 13.33, 0.03, LIME)
    add_textbox(slide, "Grab", 0.2, 0.05, 1.2, 0.35,
                font_size=16, bold=True, color=GRAB_GREEN)
    add_textbox(slide, title_right, 11.5, 0.05, 1.7, 0.35,
                font_size=13, bold=True, color=TEAL_DARK, align=PP_ALIGN.RIGHT)

def add_footer_bar(slide):
    add_rect(slide, 0, 7.2, 13.33, 0.3, LIME)

def add_section_label(slide, label):
    """Small teal pill label top-left of content area."""
    pill = add_rect(slide, 0.5, 0.6, 2.8, 0.28, TEAL_MID)
    add_textbox(slide, label, 0.55, 0.6, 2.7, 0.28,
                font_size=9, bold=True, color=WHITE)

# ── Slide builders ─────────────────────────────────────────────────────────────

def slide_cover(prs):
    layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(layout)

    # Full lime background
    add_rect(slide, 0, 0, 13.33, 7.5, LIME)

    # Decorative concentric circles (right side) – simplified rectangles approach
    cx, cy = 9.8, 3.2  # centre in inches
    for r, alpha in [(3.5, TEAL_MID), (2.6, TEAL_DARK), (1.8, TEAL_MID)]:
        shape = slide.shapes.add_shape(
            9,  # oval
            Inches(cx - r), Inches(cy - r), Inches(r * 2), Inches(r * 2)
        )
        shape.fill.background()
        shape.line.color.rgb = alpha
        shape.line.width = Pt(4)

    # Cyan accent bar (right edge)
    add_rect(slide, 12.5, 1.8, 0.83, 1.6, CYAN_LIGHT)

    # GrabHack logo text block
    add_textbox(slide, "GrabHack", 0.55, 1.6, 5, 1.1,
                font_name="Calibri", font_size=54, bold=True, color=TEAL_DARK)
    add_textbox(slide, "2.0", 4.45, 1.6, 1.4, 1.1,
                font_name="Calibri", font_size=54, bold=True, color=TEAL_MID)
    add_textbox(slide, "Shaping the Future", 0.55, 2.65, 5.5, 0.55,
                font_name="Calibri", font_size=20, bold=True, color=TEAL_DARK)

    # Divider
    add_rect(slide, 0.55, 3.35, 5.2, 0.04, TEAL_DARK)

    # Team & case study name
    add_textbox(slide, "NewAge Innovators", 0.55, 3.5, 6, 0.55,
                font_size=22, bold=True, color=TEAL_DARK)
    add_textbox(slide, "AI Saver Agent for Gig Workers", 0.55, 4.05, 7.5, 0.6,
                font_size=18, bold=False, color=TEAL_DARK)

    # Grab Family logos
    add_textbox(slide, "Grab Family", 0.55, 5.4, 2, 0.28,
                font_size=9, bold=False, color=TEAL_DARK)
    add_textbox(slide, "Grab", 0.55, 5.65, 1.0, 0.38,
                font_size=16, bold=True, color=GRAB_GREEN)
    add_textbox(slide, "GxS", 1.65, 5.65, 1.0, 0.38,
                font_size=16, bold=True, color=TEAL_DARK)

    # Platform partner
    add_textbox(slide, "Platform Partner: unstop", 10.5, 6.9, 2.7, 0.3,
                font_size=8, bold=False, color=TEAL_DARK, align=PP_ALIGN.RIGHT)

    add_footer_bar(slide)
    return slide


def slide_team(prs):
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    add_rect(slide, 0, 0, 13.33, 7.5, WHITE)
    add_header_bar(slide)
    add_footer_bar(slide)
    add_section_label(slide, "TEAM")

    add_textbox(slide, "NewAge Innovators", 0.5, 1.05, 9, 0.7,
                font_size=32, bold=True, color=TEAL_DARK)
    add_textbox(slide, "AI Saver Agent for Gig Workers", 0.5, 1.75, 10, 0.5,
                font_size=20, bold=False, color=TEAL_MID)

    # Divider
    add_rect(slide, 0.5, 2.4, 12.3, 0.05, LIME)

    members = [
        ("Member 1", "Role / Expertise"),
        ("Member 2", "Role / Expertise"),
        ("Member 3", "Role / Expertise"),
        ("Member 4", "Role / Expertise"),
    ]
    for i, (name, role) in enumerate(members):
        col = i % 2
        row = i // 2
        lx = 0.7 + col * 6.4
        ty = 2.7 + row * 1.6
        box = add_rect(slide, lx, ty, 5.8, 1.3, GRAY_BG)
        add_rect(slide, lx, ty, 0.12, 1.3, TEAL_MID)
        add_textbox(slide, name, lx + 0.3, ty + 0.15, 5.2, 0.45,
                    font_size=16, bold=True, color=TEAL_DARK)
        add_textbox(slide, role, lx + 0.3, ty + 0.6, 5.2, 0.4,
                    font_size=13, bold=False, color=BLACK)
    return slide


def slide_problem(prs):
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    add_rect(slide, 0, 0, 13.33, 7.5, WHITE)
    add_header_bar(slide)
    add_footer_bar(slide)
    add_section_label(slide, "PROBLEM STATEMENT")

    add_textbox(slide, "The Problem", 0.5, 1.05, 9, 0.6,
                font_size=28, bold=True, color=TEAL_DARK)
    add_rect(slide, 0.5, 1.7, 12.3, 0.05, LIME)

    points = [
        ("Irregular Income", "Gig workers face unpredictable earnings, making it hard to plan savings or meet financial goals."),
        ("No Safety Net", "Lack of employer-backed benefits leaves gig workers exposed to financial shocks and emergencies."),
        ("Low Financial Literacy", "Most gig workers have limited access to personalised financial guidance or smart savings tools."),
        ("Missed Micro-Savings Opportunities", "High transaction volumes present untapped opportunities for automated, low-friction savings."),
    ]
    for i, (title, body) in enumerate(points):
        col = i % 2
        row = i // 2
        lx = 0.6 + col * 6.45
        ty = 1.95 + row * 2.35
        add_rect(slide, lx, ty, 6.1, 2.1, GRAY_BG)
        add_rect(slide, lx, ty, 6.1, 0.07, TEAL_MID)
        add_textbox(slide, f"⚡  {title}", lx + 0.2, ty + 0.15, 5.7, 0.45,
                    font_size=14, bold=True, color=TEAL_DARK)
        add_textbox(slide, body, lx + 0.2, ty + 0.6, 5.7, 1.3,
                    font_size=11, bold=False, color=BLACK, wrap=True)
    return slide


def slide_solution(prs):
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    add_rect(slide, 0, 0, 13.33, 7.5, WHITE)
    add_header_bar(slide)
    add_footer_bar(slide)
    add_section_label(slide, "SOLUTION")

    add_textbox(slide, "AI Saver Agent for Gig Workers", 0.5, 1.05, 10, 0.65,
                font_size=26, bold=True, color=TEAL_DARK)
    add_rect(slide, 0.5, 1.75, 12.3, 0.05, LIME)

    add_textbox(slide,
        "An intelligent, conversational AI agent embedded in Grab SuperApp that helps gig workers "
        "automatically save money in GXS Bank accounts using personalised micro-savings rules, "
        "behavioural nudges, and real-time income analytics.",
        0.5, 1.9, 12.3, 0.9, font_size=13, color=BLACK, wrap=True)

    features = [
        ("🤖", "Smart Savings Agent", "Conversational AI understands income patterns and auto-saves based on configurable rules."),
        ("📊", "Income Analytics", "Real-time dashboard of earnings, spending trends and savings progress across Grab rides & deliveries."),
        ("🎯", "Goal-Based Savings", "Set savings goals (emergency fund, insurance, festival) with AI-driven milestones & motivation."),
        ("💡", "Behavioural Nudges", "Context-aware notifications at the right moment – post-trip, payday, or low-activity periods."),
        ("🔒", "GXS Integration", "Seamless one-tap deposit into GXS savings pots with competitive interest rates."),
        ("🏆", "Gamification", "Streaks, badges and leaderboards to drive consistent saving habits among gig workers."),
    ]
    for i, (icon, title, body) in enumerate(features):
        col = i % 3
        row = i // 3
        lx = 0.5 + col * 4.28
        ty = 3.05 + row * 1.85
        add_rect(slide, lx, ty, 4.0, 1.65, GRAY_BG)
        add_rect(slide, lx, ty, 0.1, 1.65, TEAL_MID)
        add_textbox(slide, f"{icon}  {title}", lx + 0.2, ty + 0.1, 3.7, 0.4,
                    font_size=12, bold=True, color=TEAL_DARK)
        add_textbox(slide, body, lx + 0.2, ty + 0.5, 3.7, 1.0,
                    font_size=10, color=BLACK, wrap=True)
    return slide


def slide_architecture(prs):
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    add_rect(slide, 0, 0, 13.33, 7.5, WHITE)
    add_header_bar(slide)
    add_footer_bar(slide)
    add_section_label(slide, "TECHNICAL ARCHITECTURE")

    add_textbox(slide, "How It Works", 0.5, 1.05, 9, 0.6,
                font_size=28, bold=True, color=TEAL_DARK)
    add_rect(slide, 0.5, 1.7, 12.3, 0.05, LIME)

    layers = [
        ("User Layer", "Grab SuperApp (Driver / Delivery Partner)", TEAL_MID),
        ("AI Agent Layer", "Conversational Agent · Income Analyser · Nudge Engine · Goal Tracker", TEAL_DARK),
        ("Data Layer", "Grab Trip Data · GXS Transaction API · User Profile & Preferences", CYAN_LIGHT),
        ("Banking Layer", "GXS Bank API – Savings Pots · Interest Engine · KYC / Compliance", GRAB_GREEN),
    ]
    for i, (label, desc, color) in enumerate(layers):
        ty = 1.95 + i * 1.18
        add_rect(slide, 0.5, ty, 12.3, 0.95, GRAY_BG)
        add_rect(slide, 0.5, ty, 0.15, 0.95, color)
        add_textbox(slide, label, 0.8, ty + 0.05, 2.8, 0.38,
                    font_size=12, bold=True, color=color)
        add_textbox(slide, desc, 0.8, ty + 0.45, 11.5, 0.38,
                    font_size=11, color=BLACK)
        if i < len(layers) - 1:
            add_textbox(slide, "▼", 6.4, ty + 0.95, 0.5, 0.3,
                        font_size=12, bold=True, color=TEAL_MID, align=PP_ALIGN.CENTER)

    add_textbox(slide,
        "Tech Stack: Python · FastAPI · LangChain / Claude API · PostgreSQL · Redis · React Native",
        0.5, 6.65, 12.3, 0.38, font_size=10, color=TEAL_DARK)
    return slide


def slide_impact(prs):
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    add_rect(slide, 0, 0, 13.33, 7.5, WHITE)
    add_header_bar(slide)
    add_footer_bar(slide)
    add_section_label(slide, "IMPACT & VALUE PROPOSITION")

    add_textbox(slide, "Why It Matters", 0.5, 1.05, 9, 0.6,
                font_size=28, bold=True, color=TEAL_DARK)
    add_rect(slide, 0.5, 1.7, 12.3, 0.05, LIME)

    metrics = [
        ("3M+", "Grab gig workers\nacross SEA"),
        ("40%", "Average income\nvolatility for gig workers"),
        ("8x", "Higher savings rate\nwith AI nudging (research)"),
        ("$500M+", "Potential GXS\ndeposit growth"),
    ]
    for i, (num, label) in enumerate(metrics):
        lx = 0.6 + i * 3.18
        add_rect(slide, lx, 2.0, 2.85, 2.0, TEAL_DARK)
        add_textbox(slide, num, lx + 0.1, 2.15, 2.65, 0.9,
                    font_size=34, bold=True, color=LIME, align=PP_ALIGN.CENTER)
        add_textbox(slide, label, lx + 0.1, 3.05, 2.65, 0.8,
                    font_size=11, color=WHITE, align=PP_ALIGN.CENTER, wrap=True)

    benefits = [
        ("For Gig Workers", [
            "Build financial resilience with zero friction",
            "Personalised savings plans aligned to income cycles",
            "Access to GXS high-yield savings accounts",
        ]),
        ("For Grab / GxS", [
            "Increase GXS deposit book & ARPU",
            "Deepen platform stickiness & worker loyalty",
            "ESG impact: financial inclusion for underserved segment",
        ]),
    ]
    for i, (title, pts) in enumerate(benefits):
        lx = 0.6 + i * 6.4
        add_rect(slide, lx, 4.3, 6.0, 2.55, GRAY_BG)
        add_rect(slide, lx, 4.3, 6.0, 0.08, TEAL_MID)
        add_textbox(slide, title, lx + 0.2, 4.42, 5.5, 0.4,
                    font_size=14, bold=True, color=TEAL_DARK)
        for j, pt in enumerate(pts):
            add_textbox(slide, f"•  {pt}", lx + 0.25, 4.92 + j * 0.55, 5.5, 0.45,
                        font_size=11, color=BLACK, wrap=True)
    return slide


def slide_go_to_market(prs):
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    add_rect(slide, 0, 0, 13.33, 7.5, WHITE)
    add_header_bar(slide)
    add_footer_bar(slide)
    add_section_label(slide, "GO-TO-MARKET")

    add_textbox(slide, "Roadmap & Roll-out", 0.5, 1.05, 9, 0.6,
                font_size=28, bold=True, color=TEAL_DARK)
    add_rect(slide, 0.5, 1.7, 12.3, 0.05, LIME)

    phases = [
        ("Phase 1\nMonth 1–3", "MVP Pilot",
         "• 5,000 Grab drivers in Singapore\n• Basic income analytics\n• Manual savings rules\n• GXS pot integration"),
        ("Phase 2\nMonth 4–6", "AI Enhancement",
         "• AI agent with natural language\n• Smart micro-savings automation\n• Goal-based savings pots\n• Nudge engine v1"),
        ("Phase 3\nMonth 7–12", "Scale & Expand",
         "• All Grab gig workers in SG\n• Expand to MY, ID, PH\n• Gamification & community\n• Insurance micro-products"),
    ]
    for i, (phase, title, detail) in enumerate(phases):
        lx = 0.55 + i * 4.28
        add_rect(slide, lx, 2.0, 3.95, 4.8, GRAY_BG)
        add_rect(slide, lx, 2.0, 3.95, 0.8, TEAL_DARK)
        add_textbox(slide, phase, lx + 0.15, 2.05, 3.65, 0.7,
                    font_size=11, bold=True, color=LIME, wrap=True)
        add_textbox(slide, title, lx + 0.15, 2.85, 3.65, 0.45,
                    font_size=14, bold=True, color=TEAL_DARK)
        add_textbox(slide, detail, lx + 0.15, 3.35, 3.65, 3.3,
                    font_size=10, color=BLACK, wrap=True)
        if i < len(phases) - 1:
            add_textbox(slide, "→", lx + 3.95, 4.0, 0.4, 0.5,
                        font_size=18, bold=True, color=TEAL_MID)
    return slide


def slide_closing(prs):
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    add_rect(slide, 0, 0, 13.33, 7.5, TEAL_DARK)

    # Decorative circle
    shape = slide.shapes.add_shape(9, Inches(8.5), Inches(0.5), Inches(7), Inches(7))
    shape.fill.background()
    shape.line.color.rgb = TEAL_MID
    shape.line.width = Pt(3)

    add_rect(slide, 0, 6.9, 13.33, 0.6, LIME)

    add_textbox(slide, "GrabHack 2.0", 0.6, 1.3, 8, 0.7,
                font_size=28, bold=True, color=LIME)
    add_textbox(slide, "NewAge Innovators", 0.6, 2.1, 9, 0.8,
                font_size=38, bold=True, color=WHITE)
    add_textbox(slide, "AI Saver Agent for Gig Workers", 0.6, 2.95, 9, 0.6,
                font_size=22, bold=False, color=CYAN_LIGHT)
    add_rect(slide, 0.6, 3.7, 5, 0.06, LIME)
    add_textbox(slide, "Empowering gig workers to save smarter,\nlive better, and build financial resilience.",
                0.6, 3.85, 8.5, 0.95, font_size=16, color=WHITE, wrap=True)

    add_textbox(slide, "Thank You", 0.6, 5.6, 4, 0.6,
                font_size=28, bold=True, color=LIME)
    add_textbox(slide, "Grab  |  GxS  |  unstop", 0.6, 6.9, 4, 0.45,
                font_size=11, bold=False, color=TEAL_DARK)
    return slide


# ── Build & save ───────────────────────────────────────────────────────────────
def main():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_cover(prs)
    slide_team(prs)
    slide_problem(prs)
    slide_solution(prs)
    slide_architecture(prs)
    slide_impact(prs)
    slide_go_to_market(prs)
    slide_closing(prs)

    out = "/home/user/saver/NewAge_Innovators_AI_Saver_Agent_GrabHack2.pptx"
    prs.save(out)
    print(f"Saved: {out}")


if __name__ == "__main__":
    main()
